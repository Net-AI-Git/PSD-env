import logging
from typing import List
from pptx import Presentation
from pptx.util import Inches
import os
from copy import deepcopy

# Configure logger
logger = logging.getLogger(__name__)


def _duplicate_slide(prs, slide_to_copy):
    """
    Duplicates a slide in a presentation.

    Why (Purpose and Necessity):
    The python-pptx library does not provide a direct method for duplicating
    a slide. This utility function manually creates a copy by adding a new
    slide with the same layout and then copying over all the shapes from the
    source slide, preserving their properties.

    What (Implementation Details):
    It takes a presentation object and a source slide. It adds a new blank
    slide using the source slide's layout. Then, it iterates through every
    shape in the source slide, creates a deep copy of its element, and adds
    it to the new slide's shape tree.

    Args:
        prs (pptx.presentation.Presentation): The presentation object.
        slide_to_copy (pptx.slide.Slide): The slide to be duplicated.

    Returns:
        pptx.slide.Slide: The newly created duplicate slide.
    """
    dest_slide = prs.slides.add_slide(slide_to_copy.slide_layout)

    for shape in slide_to_copy.shapes:
        new_el = deepcopy(shape.element)
        dest_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    return dest_slide


def create_presentation_from_images(image_paths: List[str], output_dir: str) -> None:
    """
    Creates a PowerPoint presentation from a list of image paths using a template.

    Why (Purpose and Necessity):
    This function gathers all generated plot images into a single PowerPoint
    file based on a pre-defined template. This provides a consistent, branded
    way to review and present the visual results of the optimization process.

    What (Implementation Details):
    The function loads a template 'format.pptx' from the 'Data' directory.
    It uses the first slide of this template as a master for duplication.
    For each image, it duplicates the master slide, extracts a title from the
    image filename, finds a placeholder shape with 'A title' and replaces it.
    It then adds the image to the duplicated slide. After all images are
    processed, the original master slide is removed from the presentation,
    which is then saved to the specified output directory. If the template is
    not found, it falls back to creating a blank presentation.

    Args:
        image_paths (List[str]): A list of absolute paths to the image files.
        output_dir (str): The directory where the final presentation
                          file will be saved.

    Returns:
        None
    """
    logger.info("Starting PowerPoint presentation generation...")
    template_path = os.path.join("Data", "format.pptx")

    try:
        # --- 1. Load Presentation (from template or blank) ---
        if os.path.exists(template_path):
            logger.info(f"Loading PowerPoint template from '{template_path}'...")
            prs = Presentation(template_path)
        else:
            logger.warning(f"Template '{template_path}' not found. Creating a new blank presentation.")
            prs = Presentation()

        if not prs.slides:
             logger.warning("Presentation template is empty. Cannot proceed.")
             # Create at least one slide to avoid errors
             prs.slides.add_slide(prs.slide_layouts[6])

        # --- 2. Prepare slides and sort images ---
        master_slide = prs.slides[0]
        details_images = sorted([p for p in image_paths if 'details' in os.path.basename(p).lower()])
        other_images = sorted([p for p in image_paths if 'details' not in os.path.basename(p).lower()])
        sorted_image_paths = other_images + details_images

        if not sorted_image_paths:
            logger.warning("No images found to add to the presentation.")
            return

        # --- 3. Process each image: duplicate, update title, add image ---
        for image_path in sorted_image_paths:
            if not os.path.exists(image_path):
                logger.warning(f"Image not found, skipping: {image_path}")
                continue

            logger.info(f"Processing image: {os.path.basename(image_path)}")
            new_slide = _duplicate_slide(prs, master_slide)

            # Extract title (first part of the filename before '.')
            base_name = os.path.basename(image_path)
            new_title = base_name.split('.')[0]

            # Find and replace the title placeholder
            for shape in new_slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip() == "A title":
                    logger.info(f"Updating title placeholder to '{new_title}'")
                    shape.text_frame.text = new_title
                    break

            # Add and center the image on the new slide
            slide_height = prs.slide_height
            slide_width = prs.slide_width
            pic = new_slide.shapes.add_picture(image_path, Inches(1), Inches(1))

            img_aspect_ratio = pic.width / pic.height
            slide_aspect_ratio = slide_width / slide_height

            if img_aspect_ratio > slide_aspect_ratio:
                new_width = slide_width * 0.9
                new_height = int(new_width / img_aspect_ratio)
            else:
                new_height = slide_height * 0.8
                new_width = int(new_height * img_aspect_ratio)

            # Ensure all position and size values are integers
            pic.left = int((slide_width - new_width) / 2)
            pic.top = int((slide_height - new_height) / 2)
            pic.width = int(new_width)
            pic.height = int(new_height)

        # --- 4. Remove the original template slide ---
        logger.info("Finished processing all images. Removing template slide.")
        # Access the slide list's XML element and remove the first slide ID
        slides = prs.slides
        if len(slides._sldIdLst) > 0:
            slide_id_to_remove = slides._sldIdLst[0]
            slides._sldIdLst.remove(slide_id_to_remove)

        # --- 5. Save the final presentation ---
        output_path = os.path.join(output_dir, "optimization_summary.pptx")
        prs.save(output_path)
        logger.info(f"Successfully saved presentation to {output_path}")

    except Exception as e:
        logger.error(f"Failed to create PowerPoint presentation. Error: {e}", exc_info=True)
